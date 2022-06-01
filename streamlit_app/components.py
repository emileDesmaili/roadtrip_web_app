#library imports


import re
import random
import pandas as pd
from dataclasses import dataclass
from typing import Optional, List
import requests
from bs4 import BeautifulSoup
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from io import BytesIO
import datetime
import numpy as np
from IPython.display import HTML
from multiprocessing.dummy import Pool as ThreadPool
from multiprocessing import Pool
from pytrends.request import TrendReq
import networkx as nx
from streamlit_agraph import agraph, Node, Edge, Config
import smtplib
from email.mime.text import MIMEText

import json
from textblob import TextBlob
from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
from dateutil import parser
from st_aggrid import AgGrid
from st_aggrid.grid_options_builder import GridOptionsBuilder
from googletrans import Translator
from sklearn.decomposition import PCA
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from geopy.geocoders import Nominatim
from geopy.extra.rate_limiter import RateLimiter
from streamlit_tags import st_tags, st_tags_sidebar
from streamlit.state.session_state import SessionState
import unidecode
from pptx import Presentation 
from pptx.util import Inches, Pt
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator
import matplotlib.pyplot as plt
from nltk.corpus import stopwords
import spacy
from spacy import displacy
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline, AutoModelForSequenceClassification
import transformers
from sentence_transformers import SentenceTransformer, util






# checkbox function
def create_checkbox(field,data):
    """creates a streamlit checkbox from a string in the streamlit dataframe

    Args:
        field (string): 

    Returns:
        streamlit object: 
    """
    field_data = data[field].unique()
    container = st.sidebar.container()
    all = st.sidebar.checkbox(f"Select all {field}s", value = True)
    if all:
        selected_options = container.multiselect(f"Select one or more {field}:",
         field_data,field_data)
    else:
        selected_options =  container.multiselect(f"Select one or more {field}:",
        field_data)
    return selected_options


def generate_shareholders_df(json_dict):
    if json_dict['CSH_TOTAL_PCT']:
        pct = json_dict['CSH_TOTAL_PCT'] 
    elif json_dict['DUO_TOTAL_PCT']: 
        pct = json_dict['DUO_TOTAL_PCT']
    else:
        pct= None
    if json_dict['CSH_ENTITY_TYPE']:
        entity_type = json_dict['CSH_ENTITY_TYPE']
    elif json_dict['DUO_ENTITY_TYPE']:
        entity_type = json_dict['DUO_ENTITY_TYPE']
    else:
        entity_type = None

    data = zip(pct,entity_type)
    df = pd.DataFrame(data).transpose()
    if json_dict['CSH_NAME']:
        df.columns = json_dict['CSH_NAME']
    elif json_dict['DUO_NAME']:
        df.columns = json_dict['DUO_NAME']

    return df

def generate_shareholders_treemap(json_dict):
    GUO = json_dict['GUO_NAME']*3 if json_dict['GUO_NAME'] else [None]*3
    DUO = json_dict['DUO_NAME']*3 if json_dict['DUO_NAME'] else [None]*3
    ISH = json_dict['ISH_NAME']*3 if json_dict['ISH_NAME'] else json_dict['DUO_NAME']*3 
    pct = [0.1,0.1,1]
    df = pd.DataFrame(dict(GUO=GUO, DUO=DUO, ISH=ISH, pct=pct))

    fig = px.treemap(df, path=['GUO', 'DUO', 'ISH'], values='pct')
    st.plotly_chart(fig)

def graph_shareholders(firm,json_dict):
        nodes = []
        edges = []

        nodes.append( Node(id=firm, label = firm, size=1000, color = "#0a2357" ))
        if json_dict['GUO_NAME']:
            for i in range(len(json_dict['GUO_NAME'])):
                name = json_dict['GUO_NAME'][i]
                try:
                    pct = "{:.0%}".format(float(json_dict['GUO_DIRECT_PCT'][i])/100)
                except:
                    pct = None
                    
                nodes.append( Node(id=str(name),
                label=name,
                size=800,
                color="#fc0349")
                ) 
                edges.append( Edge(source=str(name), 
                label=pct, 
                target=firm )
                )       
        if json_dict['DUO_NAME']:
            for i in range(len(json_dict['DUO_NAME'])):
                name = json_dict['DUO_NAME'][i]
                try:
                    
                    pct = "{:.0%}".format(float(json_dict['DUO_DIRECT_PCT'][i])/100)
                except:
                    pct = None
                nodes.append( Node(id=str(name),
                label=name,
                size=800,
                color="#fc0349")
                )
                edges.append( Edge(source=str(name), 
                label=pct, 
                target=firm )
                )
        if json_dict['ISH_NAME']:
            for i in range(len(json_dict['ISH_NAME'])):
                name = json_dict['ISH_NAME'][i]
                try:

                    pct = "{:.0%}".format(float(json_dict['ISH_DIRECT_PCT'][i])/100)
                except:
                    pct = None
                nodes.append( Node(id=str(name),
                label=name,
                size=800,
                color="#fc0349")
                )
                edges.append( Edge(source=str(name), 
                label=pct, 
                target=firm )
                )




        config = Config(width=1000, 
                        height=500, 
                        directed=True,
                        nodeHighlightBehavior=True, 
                        highlightColor="#778ebd", # or "blue"
                        collapsible=True,
                        node={'labelProperty':'label'},
                        link={'labelProperty': 'label', 'renderLabel': True}
                        # **kwargs e.g. node_size=1000 or node_color="blue"
                        ) 

        return_value = agraph(nodes=nodes, 
                            edges=edges, 
                            config=config)
        return return_value





@st.cache(show_spinner=False)
def get_linkedin(name):
    USER_AGENT = (
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10.14; rv:65.0) Gecko/20100101 Firefox/66.0"
    )
    query = f"{name}"
    query = query.replace(" ", "+")
    URL = f"https://google.com/search?q={query}+linkedin"
    headers = {"user-agent": USER_AGENT}
    
    query = requests.get(URL,headers=headers)
    
    soup = BeautifulSoup(query.content, "html.parser", from_encoding='utf8')
    if soup.find('cite'):
        links = soup.find('cite').text
        links = links.split(' › ')[1]
        linkedin='https://fr.linkedin.com/in/'+links
        return linkedin
    else:
        return 'https://fr.linkedin.com/in/'


def generate_management_df(json_dict):

    k = json_dict['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred']
    v = json_dict['CPYCONTACTS_MEMBERSHIP_Function']
    df = pd.DataFrame(k)
    if len(df.columns)>0:
        df.columns=['Name']
        df['Title'] = v
        df = df.drop_duplicates(subset=['Name'], ignore_index=True)
        df['linkedin']=df['Name'].apply(get_linkedin)
        return df
    else:
        return df

def generate_financials_df(json_dict):
    data = []
    keys = []
    columns=['Net Income', 'Turnover','EBITDA','Working Capital', 'Gearing','Leverage Ratio','Financial Debt',
    'Cash Flow','Gross Margin','EBITDA Margin','ROA','ROE','Cash','EBIT','Operating Margin','Liquidity Ratio','Taxes']
    if json_dict['Years']:
        for year in json_dict['Years']:
            keys.append(json_dict['Years'][year]['FISCAL_YEAR'])
            data.extend([[
                json_dict['Years'][year]['PL'],
                json_dict['Years'][year]['OPRE'],
                json_dict['Years'][year]['EBTA'],
                json_dict['Years'][year]['WKCA'],
                json_dict['Years'][year]['GEAR'],
                json_dict['Years'][year]['SHLQ'],
                json_dict['Years'][year]['LTDB'],
                json_dict['Years'][year]['CF'],
                json_dict['Years'][year]['GRMA'],
                json_dict['Years'][year]['ETMA'],
                json_dict['Years'][year]['ROA'],
                json_dict['Years'][year]['ROE'],
                json_dict['Years'][year]['CASH'],
                json_dict['Years'][year]['OPPL'],
                json_dict['Years'][year]['EBMA'],
                json_dict['Years'][year]['LIQR'],
                json_dict['Years'][year]['TAXA'],
                ]])        
        df = pd.DataFrame(data)
        df.index = keys
        df.columns = columns   
        return df
    else:
        return pd.DataFrame(data)



@st.cache()
def get_news(client, n):
    title = []
    url = []
    date = []
    import xmltodict
    google_actu_url=f"https://news.google.com/rss/search?q=%7B{client}%7D&hl=fr&gl=FR&ceid=FR:fr"
    google_actu_url = google_actu_url.replace(" ", "+")
    response = requests.get(google_actu_url)
    decoded_response = response.content.decode('utf-8')
    response_json = json.loads(json.dumps(xmltodict.parse(decoded_response)))
    for i in range(0,n):
        try:
            url.append(response_json['rss']['channel']['item'][i]['link'])
        except (IndexError,ValueError, KeyError):
            return 'Oops, something went wrong... try again!'
        title.append(response_json['rss']['channel']['item'][i]['title'])
        date.append(response_json['rss']['channel']['item'][i]['pubDate'])
    df = pd.DataFrame([date,title, url]).transpose()
    df[0] = df[0].apply(parser.parse)
    return df.set_index(df[0])

def st_write_news(client,n=5):
    news_df = get_news(client,n=n)
    st.markdown("""
    <style>
    .big-font {
    font-size:14px;
        }
    </style>
    """, unsafe_allow_html=True)
    if isinstance(news_df,str):
        st.write('Oopps something went wrong... try again!')
    else:
        for i in range(0,n):
            url = news_df[2][i]+' target="_blank"'
            text = news_df[1][i]
            string = f"<p class='big-font'><a href={url}>{text}</a> </p>"
            st.markdown(string, unsafe_allow_html=True)

def plot_news_sentiment(client,news_state_full):
    # analyzer = SentimentIntensityAnalyzer()
    # news_df = get_news(client, n)
    # news_df['sentimentBlob'] = news_df[1].apply(lambda news: TextBlob(news).sentiment.polarity)
    # news_df['sentimentVader'] = news_df[1].apply(lambda news: analyzer.polarity_scores(news)['compound'])
    # news_df['sentiment'] = 0.5*news_df['sentimentBlob']+0.5*news_df['sentimentVader']
    # # translator = Translator()
    # # news_df['EN'] = news_df[1].apply(lambda news: translator.translate(news).text)
    # # news_df['sentimentBlob_EN'] = news_df['EN'].apply(lambda news: TextBlob(news).sentiment.polarity)
    # # news_df['sentimentVader_EN'] = news_df['EN'].apply(lambda news: analyzer.polarity_scores(news)['compound'])
    # # news_df['sentiment_EN'] = 0.5*news_df['sentimentBlob_EN']+0.5*news_df['sentimentVader_EN']
    # senti_fig = px.line(news_df.sort_index(), x=0, y=news_df['sentiment'])
    # senti_fig.update_layout(yaxis_title=None, xaxis_title=None) 
    # st.plotly_chart(senti_fig, use_container_width=False)

    #random
    # brownian = np.cumsum(np.random.randn(200))
    # fig = px.line(pd.Series(brownian),y=0)

    split = min(1,len(client)-2)
    client=' '.join(client.split()[:split])
    pytrends = TrendReq(hl='en-US', tz=360) 
    kw_list = [client] # list of keywords to get data 
    pytrends.build_payload(kw_list, cat=0, timeframe='today 12-m')

    data = pytrends.interest_over_time() 
    data = data.reset_index() 


    fig = px.line(data, x="date", y=[client])

    fig.update_traces(line=dict(width=3))



    
    fig.update_layout(yaxis_title=None, xaxis_title=None,
    width=100, height=150) 
    fig.update_layout(margin={"r":0,"t":0,"l":0,"b":0})
    fig.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    # fig.add_annotation(text="FAKE DATA", align="center", opacity=0.20, font=dict(
    #         family="Helvetica",
    #         size=65,
    #         color="Black"
    #         ),
    #         xref="paper", yref="paper",
    #         x=0.5, y=0.8, showarrow=False)
    st.write(f'##### Search Interest Over Time for {client} from Google Trends')
    try:
        st.metric('Google Trend Score',"{:,.1f}".format(data[client].iloc[-1]),delta="{:.0%}".format((data[client].iloc[-1]/data[client].iloc[-2])-1))
    except:
        pass

    st.plotly_chart(fig, use_container_width=True)




def plot_clusters(data,dim):
    features = ['Turnover','Sector Code']
    features = ['Sector Code']
    np.random.seed(42)
    scaler = StandardScaler()
    kmeans_df = data.set_index('Name')[features].dropna()
    kmeans_df['Sector Code'] = kmeans_df['Sector Code'].str.replace('.','').str[:-1] #not great
    #kmeans_df['min'] = min([int(s) for s in data['Employee Range'][0].split() if s.isdigit()])
    #kmeans_df['max'] = max([int(s) for s in data['Employee Range'][0].split() if s.isdigit()])
    kmeans_df = pd.DataFrame(scaler.fit_transform(kmeans_df),columns=kmeans_df.columns, index=kmeans_df.index)
    k=5 # no elbow/silhouette here...
    make_cluster = KMeans(n_clusters=k)
    KM_fit = make_cluster.fit_transform(kmeans_df)
    dim_red = PCA(n_components=dim)
    km_viz = dim_red.fit_transform(KM_fit) 
    if dim == 3:
        cluster_plot_df = pd.DataFrame(np.c_[km_viz, make_cluster.labels_+1],index=kmeans_df.index,columns=['PC1','PC2','PC3','cluster'])
        fig = px.scatter_3d(cluster_plot_df, x='PC1', y='PC2',z='PC3', color='cluster', hover_name =cluster_plot_df.index,
        hover_data={"PC1":False,"PC2":False, "PC3":False}, color_discrete_sequence=px.colors.sequential.Electric)
        fig.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
        fig.update_layout(scene=dict(xaxis=dict(showticklabels=False),yaxis=dict(showticklabels=False),
        zaxis=dict(showticklabels=False),
        ))
    else:
        cluster_plot_df = pd.DataFrame(np.c_[km_viz, make_cluster.labels_+1],index=kmeans_df.index,columns=['PC1','PC2','cluster'])
        fig = px.scatter(cluster_plot_df, x='PC1', y='PC2', color='cluster', hover_name =cluster_plot_df.index,
        hover_data={"PC1":False,"PC2":False}, color_discrete_sequence=px.colors.sequential.Electric)
        fig.update_layout(yaxis_title=None, xaxis_title=None, yaxis_showticklabels=False, xaxis_showticklabels=False)
        #fig.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    groups = cluster_plot_df.groupby('cluster').groups
    cols = st.columns(k)
    for i, x in enumerate(cols):
        with x:
            st.markdown("")
            st.write(f"Cluster {i+1}")
            st.markdown("***")
            st.markdown(", ".join(groups[i+1].to_list()))
    buff, figcol, buff2 = st.columns([1,4,1])
    with figcol:
        st.plotly_chart(fig)


def plot_financials(client,json_dicts):
    client_dict = json_dicts.loc[client][0]
    fin_df = generate_financials_df(client_dict)
    fin_df = fin_df.copy()
    fin_df.index = pd.to_datetime(fin_df.index, format='%Y-%m-%d %H:%M:%S')
    session_var_name = f'last_fin{client}'
    session_state_df = st.session_state[session_var_name]
    if 'Metric' in session_state_df.columns:
        session_state_df = session_state_df.drop('Metric',axis=1)
    session_state_df = session_state_df.transpose()
    session_state_df.index = pd.to_datetime(session_state_df.index, format='%Y-%m-%d %H:%M:%S')
    
    new_fin_df = fin_df.append(session_state_df).drop_duplicates().sort_index()

    fin_chart1, fin_chart2, fin_chart3 = st.columns(3)
    margins = ['EBITDA Margin','Operating Margin']
    ROIs = ['ROA','ROE']
    profits = ['Turnover','EBITDA']
    net_incomes = ['EBIT','Net Income']

    if fin_df['Cash Flow'][0] is not None:
        last_turnover_growth = "{:.0f} M€".format(fin_df['Cash Flow'][0]/1e6)
    else: 
        last_turnover_growth = '-'
    if  fin_df['EBITDA Margin'][0] is not None:
        last_EBITDA_margin = "{:.0%}".format(fin_df['EBITDA Margin'][0]/100)
    else: 
        last_EBITDA_margin = '-'
    if  fin_df['Leverage Ratio'][0] is not None:
        last_leverage = "{:.0%}".format(fin_df['Leverage Ratio'][0])
    else: 
        last_leverage = '-'
    

    fig1 = px.bar(new_fin_df.fillna(0), x=new_fin_df.index, y=profits, barmode='group',
    width=200, height=400, color_discrete_map={
        profits[0]: 'rgb(120, 186, 240)',
        profits[1]: 'rgb(26, 118, 255)'
    })
    fig1.update_layout(margin={"r":50,"t":60,"l":50,"b":50})
    fig1.update_layout(title =f'Cash Flow :   <b>{last_turnover_growth}</b>',title_x=0.45, title_y = 0.99,xaxis_title=None,yaxis_title=None) 
    fig1.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    # fig1.update_layout(legend=dict(
    # yanchor="top",
    # y=0.99,
    # xanchor="left",
    # x=0.01
    # ))
    # text1=[fin_df['Turnover'].values,fin_df['Net Income'].values]
    # for i, t in enumerate(text1):
    #     fig1.data[i].text = t
    #     fig1.data[i].textposition = 'outside'
    #     fig1.data[i].texttemplate='%{y:.2s}'
    # fig1.update_layout(uniformtext_minsize=10, uniformtext_mode='hide')

    fig2 = px.bar(new_fin_df.fillna(0), x=new_fin_df.index, y=net_incomes,barmode='group',
    width=200, height=400, color_discrete_map={
        net_incomes[0]: 'rgb(120, 186, 240)',
        net_incomes[1]: 'rgb(26, 118, 255)'
    })
    fig2.update_layout(title =f'EBITDA Margin:   <b>{last_EBITDA_margin}</b>',title_x=0.5, title_y = 0.99, xaxis_title=None,yaxis_title=None)
    fig2.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    fig2.update_layout(margin={"r":50,"t":60,"l":50,"b":50})
    # fig2.update_layout(legend=dict(
    # yanchor="top",
    # y=0.99,
    # xanchor="left",
    # x=0.01
    # ))
    # text2=[fin_df['EBITDA'][0:1],fin_df['EBIT'].iloc[0:1]]
    # for i, t in enumerate(text1):
    #     fig2.data[i].text = t
    #     fig2.data[i].textposition = 'outside'
    #     fig2.data[i].texttemplate='%{y:.2s}'
    # fig2.update_layout(uniformtext_minsize=10, uniformtext_mode='hide')


    fig3 = px.bar(new_fin_df.fillna(0), x=new_fin_df.index, y=ROIs,barmode='group',
    width=200, height=400, color_discrete_map={
        ROIs[0]: 'rgb(120, 186, 240)',
        ROIs[1]: 'rgb(26, 118, 255)'
    })
    fig3.update_layout(title =f'Leverage:   <b>{last_leverage}</b>',title_x=0.45,title_y = 0.99, xaxis_title=None,yaxis_title=None)
    fig3.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    fig3.update_layout(margin={"r":50,"t":60,"l":50,"b":50})
    # fig3.update_layout(legend=dict(
    # yanchor="bottom",
    # y=-0.3,
    # xanchor="left",
    # x=0.01
    # ))
    # text3=[fin_df['ROA'].values,fin_df['ROE'].values]
    # for i, t in enumerate(text3):
    #     fig3.data[i].text = t
    #     fig3.data[i].textposition = 'outside'
    #     fig3.data[i].texttemplate='%{y:.2s}'
    # fig3.update_layout(uniformtext_minsize=10, uniformtext_mode='hide')

    with fin_chart1: 
        st.plotly_chart(fig1, use_container_width=True)
    with fin_chart2:
        st.plotly_chart(fig2, use_container_width=True)
    with fin_chart3: 
        st.plotly_chart(fig3, use_container_width=True)
    

def plot_shareholders(client,json_dicts):
    client_dict = json_dicts.loc[client][0]
    df = generate_shareholders_df(client_dict)
    chart = px.pie(df, values=df.iloc[0], names=df.columns, template='ggplot2',
    width=150, height=150, color_discrete_sequence=px.colors.qualitative.Safe)
    chart.update_layout(margin={"r":0,"t":0,"l":0,"b":0})
    chart.update_layout(showlegend=False)
    chart.update_layout(font={'size': 15})
    st.plotly_chart(chart, use_container_width=False)


# @st.cache(show_spinner=False)
# def firm_map(data):
#         geolocator = Nominatim(user_agent="my_app")
#         geocode = RateLimiter(geolocator.geocode, min_delay_seconds=3)

#         country='France'
#         loc_string = data["Address"][0]+","+str(int(data['Zip Code'][0]))+", "+country
#         try:
#             location = geolocator.geocode(loc_string)
#             lat = location.latitude
#             lon = location.longitude
#             map_data = pd.DataFrame({'lat': [lat], 'lon': [lon]})
#         except:
#             try:
#                 loc_string = str(int(data['Zip Code'][0]))+", "+country
#                 location = geolocator.geocode(loc_string)
#                 lat = location.latitude
#                 lon = location.longitude
#                 map_data = pd.DataFrame({'lat': [lat], 'lon': [lon]})
#             except:
#                 map_data = pd.DataFrame({'lat': [48], 'lon': [2.3]})

#         return map_data
 


def get_tags(client,data):
    data_filtered = data_filtered = data[data["Name"]==client].reset_index(drop=True)
    #key=int(data[data['Name']==client].index[0])
    keywords = st_tags(
    label='',
    text='Press enter to add more',
    value=[data_filtered['Tags'][0]],
    suggestions=['Great place to work', 'AI', 'early-stage', 
                 'Data', 'ESG'],
    maxtags = 10,key=client 
    )


def plot_spider_chart(data):
    financials=['Net Income','Turnover','EBITDA','Cash Flow', 'Debt']
    data_for_spider = data[financials].rank(pct=True).fillna(0).set_index(data['Name'])
    
    fig = go.Figure()
    buff1, col1, col2, buf2 = st.columns([1,2,2,1])
    with col1:
        Firm1 = st.selectbox('firm 1',sorted(list(data['Name'])),key='spider1')
    with col2:
        Firm2 = st.selectbox('firm 2',sorted(list(data['Name'])),key='spider2')

    if Firm1:
        r1 = data_for_spider.loc[Firm1]
    else:
        r1=[0]*len(financials)
    if Firm2:
        r2 = data_for_spider.loc[Firm2]
    else:
        r2=[0]*len(financials)

    fig.add_trace(go.Scatterpolar(
        r= r1,
        theta=financials,
        fill='toself',
        name= Firm1
    ))
    fig.add_trace(go.Scatterpolar(
        r=r2,
        theta=financials,
        fill='toself',
        name=Firm2
    ))
    fig.update_layout(
    polar=dict(
        radialaxis=dict(
        visible=True,
        range=[0, 1]
        )),
    showlegend=False
    )
    fig.update_layout(
    autosize=False,
    width=500,
    height=500,
    margin=dict(
        l=50,
        r=50,
        b=50,
        t=50,
        pad=4
    )
    )
    st.plotly_chart(fig,use_container_width=True)



def plot_main_page_financials(data_filtered):
    fin_chart1, fin_chart2, fin_chart3 = st.columns(3)
    turnovers = ['Turnover']
    ebitdas = ['EBITDA']
    net_incomes = ['Net Income']
    color = ['rgb(120, 186, 240)']

    fig1 = px.bar(data_filtered, x='Name', y='Turnover', barmode='group',color_discrete_sequence=color
    )
    fig1.update_layout(title ='<b>Turnover</b>',title_x=0.4,xaxis_title=None,yaxis_title=None) 
    fig1.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    fig1.update_layout(showlegend=False, xaxis_tickangle=-90)
    fig1.update_layout(
    autosize=False,
    width=500,
    height=450,
    margin=dict(
        l=40,
        r=40,
        b=40,
        t=40,
        pad=4
    )
    )
    fig2 = px.bar(data_filtered, x='Name', y='EBITDA', barmode='group',
    color_discrete_sequence=color)
    fig2.update_layout(title ='<b>EBITDA</b>',title_x=0.4,xaxis_title=None,yaxis_title=None) 
    fig2.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    fig2.update_layout(showlegend=False, xaxis_tickangle=-90)
    fig2.update_layout(
    autosize=False,
    width=500,
    height=450,
    margin=dict(
        l=40,
        r=40,
        b=40,
        t=40,
        pad=4
    )
    )

    fig3 = px.bar(data_filtered, x='Name', y='Net Income', barmode='group',
    color_discrete_sequence=color)
    fig3.update_layout(title ='<b>Net Income</b>',title_x=0.4,xaxis_title=None,yaxis_title=None)
    fig3.update_layout({'plot_bgcolor': 'rgba(0, 0, 0, 0)','paper_bgcolor': 'rgba(0, 0, 0, 0)',})
    fig3.update_layout(showlegend=False, xaxis_tickangle=-90)
    fig3.update_layout(
    autosize=False,
    width=500,
    height=450,
    margin=dict(
        l=40,
        r=40,
        b=40,
        t=40,
        pad=4
    )
    )

    col1,col2, col3 = st.columns(3)
    with col1:
        st.plotly_chart(fig1,use_container_width=True)
    with col2:
        st.plotly_chart(fig2,use_container_width=True)
    with col3:
        st.plotly_chart(fig3,use_container_width=True)


def similarity_matrix(data):
    NAF = data['Sector Code']
    NAF = NAF.str.replace('.','').str[:-1]
    NAF = NAF.astype(float)
    def similarity(x,y):
        return min(x,y) / max(x, y)
    NAF_mat = np.zeros((len(NAF),len(NAF)))
    for i in range(len(NAF_mat)):
        for j in range(len(NAF_mat)):
            NAF_mat[i,j] = similarity(NAF[i]/100,NAF[j]/100)
    return NAF_mat

@st.cache()
def similarity_matrix2(data,nlp):

    topics = data['Activity'].replace('[source: Bureau van Dijk]', '.').apply(nlp)
    sim_mat = np.zeros((len(topics),len(topics)))
    for i in range(len(sim_mat)):
        for j in range(len(sim_mat)):
            sim_mat[i,j] = topics[i].similarity(topics[j])

    return sim_mat

@st.cache(show_spinner=False)
def similarity_matrix3(data):
    model = SentenceTransformer('sentence-transformers/paraphrase-MiniLM-L6-v2')
    topics = np.array(data['Activity'].replace('[source: Bureau van Dijk]', '.'))
    sim_mat = np.zeros((len(topics),len(topics)))
    embs = model.encode(topics,convert_to_tensor=True)


    for i in range(len(sim_mat)):
        for j in range(len(sim_mat)):

            sim_mat[i,j] = util.cos_sim(embs[i],embs[j])

    return sim_mat

def smart_search(text,data):
    
    my_bar = st.progress(5)
    model = SentenceTransformer('sentence-transformers/paraphrase-MiniLM-L6-v2')
    topics = np.array(data['Activity'].replace('[source: Bureau van Dijk]', '.'))
    my_bar.progress(50)
    text = model.encode(text,convert_to_tensor=True)
    embs = model.encode(topics,convert_to_tensor=True)
    my_bar.progress(75)
    sim_vec = np.zeros((len(topics),1))
    
    for i in range(len(topics)):
        sim_vec[i] = util.cos_sim(text, embs[i])

    my_bar.progress(100)
    
    sim_vec = pd.DataFrame(np.array(sim_vec),index=data['Name'],columns=['score'])
    sim_vec['name'] = sim_vec.index
    my_bar.empty()
    
    return sim_vec




  

def similar_firms_graph(firm,similar_firms):
    
    nodes = []
    edges = []

    nodes.append( Node(id=firm, label = firm, size=1000, color = "#0a2357" ))

    for i in range(0,len(similar_firms)):
        idx = similar_firms.index[i].split()[:3]
        nodes.append( Node(id=str(idx),
                        label=idx,
                        size=800,
                        color="#5c8bed")
                        )
        edges.append( Edge(source=str(idx), 
            label="", 
            target=firm, 
            type="STRAIGHT") 
            )


    config = Config(width=1000, 
                    height=500, 
                    directed=False,
                    nodeHighlightBehavior=True, 
                    highlightColor="#778ebd", # or "blue"
                    collapsible=True,
                    node={'labelProperty':'label'},
                    link={'labelProperty': 'label', 'renderLabel': True}
                    # **kwargs e.g. node_size=1000 or node_color="blue"
                    ) 

    return_value = agraph(nodes=nodes, 
                        edges=edges, 
                        config=config)
    return return_value


def network_graph(firm,json_dict,similar_firms):
    
    nodes = []
    edges = []

    nodes.append( Node(id=firm, label = firm, size=1000, color = "#aaa0db"))
    nodes.append( Node(id="competitors", label="Competitors", size=1000, color = "#0a2357" ))
    #nodes.append( Node(id="managers", label="Managers", size=1000, color = "#0a2357" ))
    nodes.append( Node(id="shareholders", label="Shareholders", size=1000, color = "#0a2357"))
    edges.append( Edge(source=firm, label="", target="competitors", type="STRAIGHT"))
    edges.append( Edge(source=firm, label="", target="shareholders", type="STRAIGHT"))
    #edges.append( Edge(source=firm, label="", target="managers", type="STRAIGHT"))

    #similar firms
    for i in range(0,len(similar_firms)):
        idx = similar_firms.index[i].split()[:3]
        nodes.append( Node(id=str(idx),
                        label=idx,
                        size=800,
                        color="#5c8bed")
                        )
        edges.append( Edge(source=str(idx), 
            label="", 
            target="competitors", 
            type="STRAIGHT") 
            )
    #shareholders

    if json_dict['GUO_NAME']:
        for i in range(len(json_dict['GUO_NAME'])):
            name = json_dict['GUO_NAME'][i]
            try:
                pct = "{:.0%}".format(float(json_dict['GUO_DIRECT_PCT'][i])/100)
            except:
                pct = None
                
            nodes.append( Node(id=str(name),
            label=name,
            size=800,
            color="#fc0349")
            ) 
            edges.append( Edge(source=str(name), 
            label=pct, 
            target="shareholders" )
            )       
    if json_dict['DUO_NAME']:
        for i in range(len(json_dict['DUO_NAME'])):
            name = json_dict['DUO_NAME'][i]
            try:
                
                pct = "{:.0%}".format(float(json_dict['DUO_DIRECT_PCT'][i])/100)
            except:
                pct = None
            nodes.append( Node(id=str(name),
            label=name,
            size=800,
            color="#fc0349")
            )
            edges.append( Edge(source=str(name), 
            label=pct, 
            target="shareholders" )
            )
    if json_dict['ISH_NAME']:
        for i in range(len(json_dict['ISH_NAME'])):
            name = json_dict['ISH_NAME'][i]
            try:

                pct = "{:.0%}".format(float(json_dict['ISH_DIRECT_PCT'][i])/100)
            except:
                pct = None
            nodes.append( Node(id=str(name),
            label=name,
            size=800,
            color="#fc0349")
            )
            edges.append( Edge(source=str(name), 
            label=pct, 
            target="shareholders" )
            )




    config = Config(width=1000, 
                    height=400, 
                    directed=False,
                    nodeHighlightBehavior=True, 
                    highlightColor="#778ebd", # or "blue"
                    collapsible=True,
                    node={'labelProperty':'label'},
                    link={'labelProperty': 'label', 'renderLabel': True}
                    # **kwargs e.g. node_size=1000 or node_color="blue"
                    ) 

    return_value = agraph(nodes=nodes, 
                        edges=edges, 
                        config=config)

    return return_value
    




def create_tagbox(tag_state,data):
    """creates a streamlit box from tags in a dataframe

    Args:
        tag_state (string): SessionState variable that is firm specific

    Returns:
        streamlit object: 
    """
    field_data = np.append(data['Tags'].unique(),[['Great Place to Work','Tech','AI','ESG','Healthcare']])
    field_data = set(np.append(field_data,[st.session_state[tag_state]]))
    container = st.container()
    all = st.checkbox(f"Select all tags", value = True)
    if all:
        selected_options = container.multiselect(f"Select one or more Tag:",
         field_data,field_data)
    else:
        selected_options =  container.multiselect(f"Select one or more Tag:",
        field_data)
    return selected_options


def json_to_df(json_dict):
    data = []
    def employee_range(employee_num):
        if employee_num <= 49:
            return '1-49'
        elif employee_num <= 99:
            return '50-99'
        elif employee_num <= 249:
            return '100-249'
        elif employee_num <= 499:
            return '250-499'
        elif employee_num <= 999:
            return '500-999'
        elif employee_num <= 1999:
            return '1000-1999'
        else:
            return '2000+'   
    keys = [
        'Name','Website','LEI','National ID','Sector Code','Sector description','Activity','Founded in','Employees','Address','Lat','Lon','Zip Code','City','Country',
            'Management','Shareholders']
    data.append(json_dict['NAME'].title())
    if json_dict['WEBSITE']:
        data.append('https://' + json_dict['WEBSITE'][0])
    else: 
        data.append('-')
    data.append(json_dict['LEI'])
    if json_dict['NATIONAL_ID']:
        data.append(json_dict['NATIONAL_ID'][0]+' ('+json_dict['NATIONAL_ID_LABEL'][0] +')')
    else:
        data.append('No national ID available')
    data.append(json_dict['NACE2_CORE_CODE'])
    data.append(json_dict['NACE2_CORE_LABEL'])
    data.append('\n '.join(filter(None, [json_dict['TRADE_DESCRIPTION_EN'],json_dict['PRODUCTS_SERVICES']])))
    date = json_dict['INCORPORATION_DATE']
    if date:
        data.append(datetime.datetime.strptime(date, '%Y-%m-%dT%H:%M:%S').year)
    else:
        data.append(1900)
    if json_dict['Years']['LY']['EMPL']:
        data.append(employee_range(json_dict['Years']['LY']['EMPL']))
    else:
        data.append('1-49')
    data.append(','.join(filter(None, [json_dict['ADDRESS_LINE1'],json_dict['ADDRESS_LINE2']])))
    data.append(json_dict['LATITUDE'])
    data.append(json_dict['LONGITUDE'])
    data.append(json_dict['POSTCODE'])
    data.append(json_dict['CITY'])
    data.append(json_dict['COUNTRY'])
    try:
        data.append(json_dict['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'][0] + ': ' + json_dict['CPYCONTACTS_MEMBERSHIP_Function'][0].title())
    except Exception:
        data.append('No Management Data available on Orbis')
    
    if json_dict['CSH_NAME'] is not None:
        shareholder_list = []
        
        for i in range(len(json_dict['CSH_NAME'])):
            try:
                shareholder_list.append(json_dict['CSH_NAME'][i].title() + ': ' + "{:.0%}".format(float(json_dict['CSH_TOTAL_PCT'][i])/100))
            except ValueError:
                shareholder_list.append(json_dict['CSH_NAME'][i].title() + ': ' + '-')
        data.append( ", ". join(shareholder_list))
    else:
        data.append('No Shareholder Data available on Orbis')
    

    if json_dict['Years']:
        keys.extend(['Year Financials','Turnover','Net Income','EBITDA','Cash Flow','Debt'])
        data.append(json_dict['Years']['LY']['FISCAL_YEAR'][0:4])
        data.append(json_dict['Years']['LY']['OPRE'])
        data.append(json_dict['Years']['LY']['PL'])
        data.append(json_dict['Years']['LY']['EBTA'])
        data.append(json_dict['Years']['LY']['CF'])
        data.append(json_dict['Years']['LY']['LTDB'])
    df = pd.DataFrame(data).transpose()
    df.columns = keys
    return df

@st.cache
def get_main_df(json_dicts):
    data = []
    for json_dict in json_dicts[0]:
        columns = json_to_df(json_dict).columns
        listed = json_to_df(json_dict).values.tolist()
        data.extend(listed)
    df = pd.DataFrame(data, columns = columns)

    if 'Tags' not in df:
        df['Tags'] = 'No Tag'
    return df

@st.cache(show_spinner=False)
def plot_wordcloud(topic,data):
    USER_AGENT = ("Mozilla/5.0 (Macintosh; Intel Mac OS X 10.14; rv:65.0) Gecko/20100101 Firefox/66.0")
    my_stops = ['Accueil', 'Ligne','contacter','Linkedin','Recherche','wikipedia','avis','annonce','marque','annonces','hiver','automne','officiel','site',
    'acheter','comme','com','description','rejoindre','Wikipédia','quels','bienvenue','chez','Paris','facebook','twitter','Recherches','abandonne','images',
    'Société','bilan','infogreffe','kbis','siret','siren','chiffre',"d'affaires",'tva','societe','ca','bilans','carrieres','adresse','établissements','entreprise','résultat',
    'toute','fr','contact','infos','ouverture','fermeture','homme','femme','groupe','abonnement','bon','code promo','promo','code','recrutement','général','Téléphone',
    'horaires','promotion',"c'bon",'corporation','company','corp','ltd','limited','email']
    data_filtered = data[data["Name"]==topic].reset_index(drop=True)
    address = str(data_filtered['City'][0])
    final_stopwords_list = stopwords.words('english') + stopwords.words('french') + topic.split() + my_stops + address.replace('-', ' ').split()
    topic=topic.replace(' ','+')
    headers = {"user-agent": USER_AGENT}
    url =f"https://google.com/search?q={topic}&num=50"  
    response = requests.get(url, headers= headers)
    soup = BeautifulSoup(response.content, 'html.parser')
    results = soup.find_all('h3')
    descriptions = []
    for result in results:
        try:
            description = result.get_text()
            if description != '': 
                descriptions.append(description)
        except:
            continue
    text = ' '.join(descriptions)
    if len(text)!=0:
        wordcloud = WordCloud(stopwords=set(final_stopwords_list),background_color='white',height=350, font_path='arial', colormap = 'Set2').generate(text)

        return wordcloud.to_array()

def fill_slide(client,json_dicts, slide):
        client_dict = json_dicts.loc[client][0]

        #get placeholders to write in
        a,b,c,d,e,f,shareholders,activity,description,g, m_and_a,financials,h,management, name = slide.placeholders
        name.text = client
    

        #get the data from json
        try:
            national_id = str(client_dict['NATIONAL_ID'][0]+' ('+client_dict['NATIONAL_ID_LABEL'][0] +')') if client_dict['NATIONAL_ID'] is not None else '-'
        except:
            national_id = '-'
        sector = str(client_dict['NACE2_CORE_LABEL']) + ' (' +str(client_dict['NACE2_CORE_CODE']) + ')'
        address =  ','.join(filter(None, [client_dict['ADDRESS_LINE1'],client_dict['ADDRESS_LINE2']])) + ', ' + client_dict['POSTCODE'] + ', ' + client_dict['CITY']
        try:
            creation_date = datetime.datetime.strptime(client_dict['INCORPORATION_DATE'], '%Y-%m-%dT%H:%M:%S').year
        except:
            creation_date = '-'
        if client_dict['Years']:
            employees = int(client_dict['Years']['LY']['EMPL']) if client_dict['Years']['LY']['EMPL'] is not None else '-'
            turnover = "{:.0f} M€".format(client_dict['Years']['LY']['OPRE']/1e6) if client_dict['Years']['LY']['OPRE'] is not None else '-'
            ebitda = "{:.0f} M€".format(client_dict['Years']['LY']['EBTA']/1e6) if client_dict['Years']['LY']['EBTA'] is not None else '-'
            net_income = "{:.0f} M€".format(client_dict['Years']['LY']['PL']/1e6) if client_dict['Years']['LY']['PL'] is not None else '-'
            year_financials = str(client_dict['Years']['LY']['FISCAL_YEAR'][0:4]) if client_dict['Years']['LY']['FISCAL_YEAR'] is not None else '-'
        
        #first value in list to get the "head" manager
        manager_string = client_dict['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'][0] + ': '+ client_dict['CPYCONTACTS_MEMBERSHIP_Function'][0] \
        if client_dict['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'] is not None else 'No Management Data available'

        # write it in placeholders
        activity.text = f'Nation ID: {national_id}\nSector: {sector}\nAddress: {address}\nCreation: {creation_date}\nEmployees: {employees}'
        description.text = '\n '.join(filter(None, [client_dict['TRADE_DESCRIPTION_EN'],client_dict['PRODUCTS_SERVICES']]))
        financials.text = f'Year Financials: {year_financials}\nTurnover: {turnover}\nEBITDA: {ebitda}\nNet Income: {net_income} '
        management.text =  manager_string
        
        #special case for shareholders, iterate and write directly
        shareholder_list = []
        try:
            for i in range(len(client_dict['CSH_NAME'])):
                try:
                    shareholder_string = client_dict['CSH_NAME'][i].title() + ': ' + "{:.0%}".format(float(client_dict['CSH_TOTAL_PCT'][i])/100)
                    shareholder_list.append(shareholder_string)
                except ValueError:
                    shareholder_string = client_dict['CSH_NAME'][i].title() + ': ' + "-"
                    shareholder_list.append(shareholder_string)

                joined_shareholders= "\n ". join(shareholder_list)
                shareholders.text = joined_shareholders
        except Exception:
            shareholders.text = 'No Shareholders data'





def create_ppt(clients,json_dicts):
    """creates a formatted ppt from a list of clients 

    Args:
        clients ([string]): client list
        json_dicts ([json]): master pickle containing jsons for all firms

    Returns:
        [pptx object]: [ppt presentation]
    """
    
    # Creating Object
    ppt = Presentation('streamlit_app/template_ppt.pptx') 
    
    for client1, client2 in zip(*[iter(clients)]*2):

        client_dict1 = json_dicts.loc[client1][0]
        client_dict2 = json_dicts.loc[client2][0]
        title_slidelayout = ppt.slide_master.slide_layouts[1]
        slide = ppt.slides.add_slide(title_slidelayout)
        #get placeholders to write in
        a,b,c,d,e,f,shareholders1,activity1,description1,g, m_and_a1,financials1,h,management1,i,j,k,l,m,shareholders2,activity2,description2,n, m_and_a2,financials2,o,management2,name1, name2 = slide.placeholders
        name1.text = client1
        name2.text = client2

         #get the data from json
        try:
            national_id = str(client_dict1['NATIONAL_ID'][0]+' ('+client_dict1['NATIONAL_ID_LABEL'][0] +')') 
        except:
            national_id = '-'
        sector = str(client_dict1['NACE2_CORE_LABEL']) + ' (' +str(client_dict1['NACE2_CORE_CODE']) + ')'
        address =  ','.join(filter(None, [client_dict1['ADDRESS_LINE1'],client_dict1['ADDRESS_LINE2']]))  + ', ' + client_dict1['POSTCODE'] + ', ' + client_dict1['CITY']
        try:
            creation_date = datetime.datetime.strptime(client_dict1['INCORPORATION_DATE'], '%Y-%m-%dT%H:%M:%S').year
        except:
            creation_date = '-'
        if client_dict1['Years']:
            employees = int(client_dict1['Years']['LY']['EMPL']) if client_dict1['Years']['LY']['EMPL'] is not None else '-'
            turnover = "{:.0f} M€".format(client_dict1['Years']['LY']['OPRE']/1e6) if client_dict1['Years']['LY']['OPRE'] is not None else '-'
            ebitda = "{:.0f} M€".format(client_dict1['Years']['LY']['EBTA']/1e6) if client_dict1['Years']['LY']['EBTA'] is not None else '-'
            net_income = "{:.0f} M€".format(client_dict1['Years']['LY']['PL']/1e6) if client_dict1['Years']['LY']['PL'] is not None else '-'
            year_financials = str(client_dict1['Years']['LY']['FISCAL_YEAR'][0:4]) if client_dict1['Years']['LY']['FISCAL_YEAR'] is not None else '-'
        
        #first value in list to get the "head" manager
        manager_string = client_dict1['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'][0] + ': '+ client_dict1['CPYCONTACTS_MEMBERSHIP_Function'][0] \
        if client_dict1['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'] is not None else 'No Management Data available'

        # write it in placeholders
        activity1.text = f'Nation ID: {national_id}\nSector: {sector}\nAddress: {address}\nCreation: {creation_date}\nEmployees: {employees}'
        description1.text = '\n '.join(filter(None, [client_dict1['TRADE_DESCRIPTION_EN'],client_dict1['PRODUCTS_SERVICES']]))
        financials1.text = f'Year Financials: {year_financials}\nTurnover: {turnover}\nEBITDA: {ebitda}\nNet Income: {net_income} '
        management1.text =  manager_string
        
        #special case for shareholders, iterate and write directly
        shareholder_list = []
        try:
            for i in range(len(client_dict2['CSH_NAME'])):
                try:
                    shareholder_string = client_dict2['CSH_NAME'][i].title() + ': ' + "{:.0%}".format(float(client_dict2['CSH_TOTAL_PCT'][i])/100)
                    shareholder_list.append(shareholder_string)
                except ValueError:
                    shareholder_string = client_dict2['CSH_NAME'][i].title() + ': ' + "-"
                    shareholder_list.append(shareholder_string)

                joined_shareholders= "\n ". join(shareholder_list)
                shareholders1.text = joined_shareholders
        except Exception:
            shareholders1.text = 'No Shareholders data'


        #get the data from json2
        try:
            national_id = str(client_dict2['NATIONAL_ID'][0]+' ('+client_dict2['NATIONAL_ID_LABEL'][0] +')')
        except:
            national_id = '-'
        sector = str(client_dict2['NACE2_CORE_LABEL']) + ' (' +str(client_dict2['NACE2_CORE_CODE']) + ')'
        address =  ','.join(filter(None, [client_dict2['ADDRESS_LINE1'],client_dict2['ADDRESS_LINE2']])) + ', ' + client_dict2['POSTCODE'] + ', ' + client_dict2['CITY']
        try:
            creation_date = datetime.datetime.strptime(client_dict2['INCORPORATION_DATE'], '%Y-%m-%dT%H:%M:%S').year
        except:
            creation_date='-'
        
        if client_dict2['Years']:
            employees = int(client_dict2['Years']['LY']['EMPL']) if client_dict2['Years']['LY']['EMPL'] is not None else '-'
            turnover = "{:.0f} M€".format(client_dict2['Years']['LY']['OPRE']/1e6) if client_dict2['Years']['LY']['OPRE'] is not None else '-'
            ebitda = "{:.0f} M€".format(client_dict2['Years']['LY']['EBTA']/1e6) if client_dict2['Years']['LY']['EBTA'] is not None else '-'
            net_income = "{:.0f} M€".format(client_dict2['Years']['LY']['PL']/1e6) if client_dict2['Years']['LY']['PL'] is not None else '-'
            year_financials = str(client_dict2['Years']['LY']['FISCAL_YEAR'][0:4]) if client_dict2['Years']['LY']['FISCAL_YEAR'] is not None else '-'
        
        #first value in list to get the "head" manager
        manager_string = client_dict2['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'][0] + ': '+ client_dict2['CPYCONTACTS_MEMBERSHIP_Function'][0] \
        if client_dict2['CPYCONTACTS_HEADER_FullNameOriginalLanguagePreferred'] is not None else 'No Management Data available'

        # write it in placeholders
        activity2.text = f'Nation ID: {national_id}\nSector: {sector}\nAddress: {address}\nCreation: {creation_date}\nEmployees: {employees}'
        description2.text = '\n '.join(filter(None, [client_dict2['TRADE_DESCRIPTION_EN'],client_dict2['PRODUCTS_SERVICES']]))
        financials2.text = f'Year Financials: {year_financials}\nTurnover: {turnover}\nEBITDA: {ebitda}\nNet Income: {net_income} '
        management2.text =  manager_string
        
        #special case for shareholders, iterate and write directly
        shareholder_list = []
        try:
            for i in range(len(client_dict2['CSH_NAME'])):
                try:
                    shareholder_string = client_dict2['CSH_NAME'][i].title() + ': ' + "{:.0%}".format(float(client_dict2['CSH_TOTAL_PCT'][i])/100)
                    shareholder_list.append(shareholder_string)
                except ValueError:
                    shareholder_string = client_dict2['CSH_NAME'][i].title() + ': ' + "-"
                    shareholder_list.append(shareholder_string)

                joined_shareholders= "\n ". join(shareholder_list)
                shareholders2.text = joined_shareholders
        except Exception:
            shareholders2.text = 'No Shareholders data'
    
    if len(clients)%2 !=0:
        title_slidelayout = ppt.slide_master.slide_layouts[2]
        slide = ppt.slides.add_slide(title_slidelayout)
        fill_slide(clients[-1],json_dicts,slide)

    return ppt
    
    

def color_import(val):
    color = '#e6ffe6' if val else '#ffe6e6'
    return f'background-color: {color}'

def tabs(default_tabs = [] ,default_active_tab=0):
        if not default_tabs:
            return None
        active_tab = st.radio("", default_tabs, index=default_active_tab, key='tabs')
        child = default_tabs.index(active_tab)+1
        st.markdown("""  
            <style type="text/css">
            div[role=radiogroup] > label > div:first-of-type, .stRadio > label {
               display: none;               
            }
            div[role=radiogroup] {
                flex-direction: unset
            }
            div[role=radiogroup] label {             
                border: 1px solid #999;
                background: #F8F8FF;
                padding: 4px 12px;
                border-radius: 2px 4px 0 0;
                position: relative;
                top: 1px;
                }
            div[role=radiogroup] label:nth-child(""" + str(child) + """) {    
                background: #FFF !important;
                border-bottom: 1px solid transparent;
            }            
            </style>
        """,unsafe_allow_html=True)        
        return active_tab


def get_sentiment(text,lang):
    analyzer = SentimentIntensityAnalyzer()
    textblob_score = TextBlob(text).sentiment.polarity
    vader_score = analyzer.polarity_scores(text)['compound']
    if textblob_score >0.2 and vader_score >0.2 and lang =='en':
        sentiment = "😀 Positive sentiment"
    elif textblob_score <0 or vader_score <0 and lang =='en':
        sentiment = "☹️ Negative sentiment"
    else:
        sentiment = ""
    return sentiment


def get_sentiment2(text,lang,tokenizer_sentiment,model_sentiment):
    classifier = pipeline("sentiment-analysis", model=model_sentiment, tokenizer=tokenizer_sentiment)
    try:
        sentiment = classifier(text)
        if lang =='en' and sentiment[0]['label'] =='positive':
            sentiment = '😀 Positive sentiment'
        elif lang =='en' and sentiment[0]['label'] =='negative':
            sentiment = '☹️ Negative sentiment'
        else:
            sentiment = ""
    except:
        sentiment = ""
    return sentiment

# @st.cache(hash_funcs={AutoModelForSeq2SeqLM: hash, transformers.models.bart.tokenization_bart_fast.BartTokenizerFast: hash,
# transformers.models.bart.modeling_bart.BartLearnedPositionalEmbedding : hash, AutoModelForSequenceClassification: hash, 
# transformers.models.bert.tokenization_bert_fast.BertTokenizerFast: hash, transformers.models.roberta.tokenization_roberta_fast.RobertaTokenizerFast: hash},
# suppress_st_warning=True)
def process_news(search_bar, news_dict, tokenizer_sentiment, model_sentiment):
    # first add tags like sentiment and subject codes
    news_dict_firm = news_dict[search_bar]
    sorted_dates = sorted(list(news_dict_firm.keys()), reverse=True)[:50]
    

    for date in sorted_dates:
        lang = news_dict_firm[date]['language_code']
        timestamp = news_dict_firm[date]['publication_date'][:10]
        news_dict_firm[date]['publication_datetime'] = datetime.datetime.fromtimestamp(int(timestamp))
        try:
            snippet = news_dict_firm[date]['snippet']
        except:
            snippet=None
        text = news_dict_firm[date]['body']
        article_text = text


        news_dict_firm[date]['sentiment'] = get_sentiment2(snippet,lang,tokenizer_sentiment=tokenizer_sentiment,model_sentiment=model_sentiment)
        #news_dict_firm[date]['sentiment'] = get_sentiment(snippet,lang)

            
        news_dict_firm[date]['tags'] = []
        if 'c17' in news_dict_firm[date]['subject_codes']:
            news_dict_firm[date]['tags'].append('💰 Funding')
        if 'c18' in news_dict_firm[date]['subject_codes']:
            news_dict_firm[date]['tags'].append('🏛 Ownership Changes')
        if 'cactio' in news_dict_firm[date]['subject_codes']:
            news_dict_firm[date]['tags'].append('🏦 Corporate Actions')
        if 'c411' in news_dict_firm[date]['subject_codes']:
            news_dict_firm[date]['tags'].append('👨‍💼 Management Changes')
        news_dict_firm[date]['tags'] = ','.join(news_dict_firm[date]['tags'])
    df_news = pd.DataFrame(news_dict_firm).transpose()
        
    return df_news

def filter_news(df_news,filter_lang,filter_date,filter_sentiment,filter_tag):
    
    df = df_news
    try:
        filter_lang = list(pd.Series(filter_lang).str.replace("French","fr").str.replace("English","en"))
    except:
        pass

    
    df_filtered =  df.loc[(df['language_code']).isin(filter_lang if filter_lang else df['language_code'].unique())
                        & (df['sentiment']).isin(filter_sentiment if filter_sentiment else df['sentiment'].unique())
                        & (df['tags'].str.contains('|'.join(filter_tag)).any(level=0))
                        ]
    return df_filtered


def display_news(df_news,tokenizer, model,nlp):
    WHITESPACE_HANDLER = lambda k: re.sub('\s+', ' ', re.sub('\n+', ' ', k.strip()))
    i=0
    dates = sorted(list(df_news.index),reverse=True)
    dates = dates[:10]
    for date in dates:
        
        st.container()
        title = df_news.loc[date]['title']
        st.write(f'#### {title}')
        #generic variable like language, source, date, body, snippet
        lang = df_news.loc[date]['language_code']
        source = df_news.loc[date]['source_name']
        try:
            snippet = df_news.loc[date]['snippet']
        except:
            snippet=None
        text = df_news.loc[date]['body']
        article_text = text
        time = df_news.loc[date]['publication_datetime']
        

        #list of tags like "Funding, M&A" and sentiment
        news_tags=[]
        for tag in df_news.loc[date]['tags'].split(','):
            news_tags.append(tag)
        news_tags.append(df_news.loc[date]['sentiment'])
        
        news_tags = [str(tag) for tag in news_tags] 
        news_tags_clean = []
        for tag in news_tags:
            if tag != "nan" and tag != "":
                news_tags_clean.append(tag)
        
        
        st.write(f"**{time}  -   {source}**")
        #create columns for tags, the "+3" is for design only

        cols = st.columns(len(news_tags)+3)
        j=0
        for a, x in enumerate(cols):
            with x:
                if j<len(news_tags_clean):
                    if news_tags_clean[j]=='😀 Positive sentiment':
                        st.success(news_tags_clean[j])
                    elif news_tags[j]=='☹️ Negative sentiment':
                        st.error(news_tags_clean[j])
                    else:
                        st.info(news_tags_clean[j])
            j+=1
        
        if snippet:
            if st.checkbox("AutoTag Snippet",key=i):
                #nlp = spacy.load("en_core_web_sm")
                doc1 = nlp(snippet)
                html = displacy.render(doc1, style="ent", page=False)
                st.write(html,unsafe_allow_html=True)
            else:
                st.markdown(snippet.replace('*',' '))
            
    
        with st.expander("See Full Text"):
            st.markdown(article_text)

        if st.button('🧬 Summarize Full Text', key=i):
            with st.spinner('Reading Article...'):     
                # if 'tokenizer' not in st.session_state:
                #     st.session_state['tokenizer'] = AutoTokenizer.from_pretrained("sshleifer/distilbart-cnn-12-6")
                # if 'model' not in st.session_state:
                #     st.session_state['model'] = AutoModelForSeq2SeqLM.from_pretrained("sshleifer/distilbart-cnn-12-6")
                input_ids = tokenizer(
                    [article_text],
                    return_tensors="pt",
                    padding="max_length",
                    truncation=True,
                    max_length=800
                )["input_ids"]

                output_ids = model.generate(
                    input_ids=input_ids,
                    max_length=120,
                    no_repeat_ngram_size=2,
                    num_beams=4
                )[0]

                summary = tokenizer.decode(
                    output_ids,
                    skip_special_tokens=True,
                    clean_up_tokenization_spaces=False
                )
            st.write("Summary:")
            st.info(summary)
        i+=1

# @st.cache(show_spinner=False,allow_output_mutation=True,
# hash_funcs={spacy.vocab.Vocab: hash, spacy.pipeline.tok2vec.Tok2Vec: hash, spacy.pipeline.tagger.Tagger:hash, spacy.pipeline.dep_parser.DependencyParser: hash,
# spacy.pipeline.senter.SentenceRecognizer: hash, spacy.pipeline.ner.EntityRecognizer: hash})
def NER_news(firm,news_df,nlp):
    firm = firm.title().split()[0]
    names = []
    labels = []
    #news_df = news_df[news_df['language_code']=='en']
    for doc in nlp.pipe(news_df['snippet']):
        for ent in doc.ents:
            names.append(ent.text)
            labels.append(ent.label_)
    ents_df = pd.DataFrame([names,labels]).transpose()
    ents_df.columns=['Name','Label']
    
    ents_df_persons = ents_df[ents_df['Label']=='PERSON'].groupby(['Name']).count().reset_index()
    ents_df_persons.columns=['Name','Count']
    
    ents_df_orgs = ents_df[ents_df['Label']=='ORG'].groupby(['Name']).count().reset_index()
    ents_df_orgs.columns=['Name','Count']

    ents_df_persons = ents_df_persons.loc[(ents_df_persons['Name']==ents_df_persons['Name'].str.title()) &
                                         (~ents_df_persons['Name'].str.contains(firm)) & 
                                         (~ents_df_persons['Name'].str.contains ("[^ a-zA-Z0-9]"))]
    ents_df_orgs = ents_df_orgs.loc[(ents_df_orgs['Name']==ents_df_orgs['Name'].str.title()) &  
                                    (~ents_df_orgs['Name'].str.contains(firm)) &
                                     (~ents_df_orgs['Name'].str.contains ("[^ a-zA-Z0-9]"))]


    return  ents_df_persons, ents_df_orgs

# Import smtplib for the actual sending function
def send_mail(textfile, me = 'emile.esmaili@ekimetrics.com', you='emile.esmaili@ekimetrics.com' ):



    # Open a plain text file for reading.  For this example, assume that
    # the text file contains only ASCII characters.
    # with open(textfile, 'rb') as fp:
    #     # Create a text/plain message
    #     msg = MIMEText(fp.read())

    # me == the sender's email address
    # you == the recipient's email address
    msg = pd.DataFrame([])
    msg['Subject'] = 'The contents of %s' % textfile
    msg['From'] = me
    msg['To'] = you

    # Send the message via our own SMTP server, but don't include the
    # envelope header.
    s = smtplib.SMTP('localhost')
    s.sendmail(me, [you], msg.as_string())
    s.quit()